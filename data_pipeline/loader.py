"""Document loading and optional image extraction for CourseMate.

This module turns raw course files (PDF/PPTX/DOCX/TXT) into structured
documents ready for splitting/embedding. Image extraction is performed for
PDF and PPTX slides to support multimodal retrieval.
"""
from __future__ import annotations

import os
from typing import Dict, List, Optional

import base64
import docx2txt
import fitz  # PyMuPDF
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class DocumentLoader:
    def __init__(
        self,
        data_dir: str,
        images_dir: str,
        extract_images: bool = True,
        supported_exts: Optional[List[str]] = None,
    ) -> None:
        self.data_dir = data_dir
        self.images_dir = images_dir
        self.project_root = os.path.abspath(os.path.join(self.data_dir, os.pardir))
        self.extract_images = extract_images
        self.supported_exts = supported_exts or [".pdf", ".pptx", ".docx", ".txt"]
        os.makedirs(self.images_dir, exist_ok=True)

    # ---------------------------------------------------------------------
    # Public APIs
    # ---------------------------------------------------------------------
    def load_all_documents(self) -> List[Dict]:
        """Walk data_dir and load every supported file into page/slide chunks."""
        if not os.path.exists(self.data_dir):
            print(f"[loader] data dir not found: {self.data_dir}")
            return []

        documents: List[Dict] = []

        for root, _, files in os.walk(self.data_dir):
            relative = os.path.relpath(root, self.data_dir)
            course_name = relative.split(os.sep)[0] if relative != "." else "default"

            for fname in files:
                ext = os.path.splitext(fname)[1].lower()
                if ext not in self.supported_exts:
                    continue

                fpath = os.path.join(root, fname)
                print(f"[loader] loading ({course_name}): {fpath}")
                docs = self.load_document(fpath)

                for doc in docs:
                    doc["course_name"] = course_name
                documents.extend(docs)

        return documents

    def load_document(self, file_path: str) -> List[Dict]:
        """Dispatch loader based on extension; return list of chunk dicts."""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self._load_pdf(file_path)
        if ext == ".pptx":
            return self._load_pptx(file_path)
        if ext == ".docx":
            return self._wrap_single_text_chunk(file_path, self._load_docx(file_path))
        if ext == ".txt":
            return self._wrap_single_text_chunk(file_path, self._load_txt(file_path))

        print(f"[loader] unsupported extension skipped: {file_path}")
        return []

    # ------------------------------------------------------------------
    # Per-format loaders
    # ------------------------------------------------------------------
    def _load_pdf(self, file_path: str) -> List[Dict]:
        reader = PdfReader(file_path)
        pdf_name = os.path.splitext(os.path.basename(file_path))[0]

        fitz_doc = None
        if self.extract_images:
            try:
                fitz_doc = fitz.open(file_path)
            except Exception as exc:
                print(f"[loader] fitz open failed: {exc}")

        chunks: List[Dict] = []
        for page_idx, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            images = (
                self._extract_images_from_pdf_page(fitz_doc, page_idx - 1, pdf_name)
                if fitz_doc
                else []
            )

            content = f"--- Page {page_idx} ---\n{text}\n"
            if images:
                content += f"\n[This page contains {len(images)} related image(s)]\n"

            chunks.append(
                {
                    "content": content,
                    "filename": os.path.basename(file_path),
                    "filepath": file_path,
                    "filetype": ".pdf",
                    "page_number": page_idx,
                    "images": images,
                }
            )

        if fitz_doc:
            fitz_doc.close()
        return chunks

    def _load_pptx(self, file_path: str) -> List[Dict]:
        prs = Presentation(file_path)
        pptx_name = os.path.splitext(os.path.basename(file_path))[0]
        chunks: List[Dict] = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            text_parts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
            text = "\n".join(text_parts)

            images = (
                self._extract_images_from_pptx_slide(slide, slide_idx, pptx_name)
                if self.extract_images
                else []
            )

            content = f"--- Slide {slide_idx} ---\n{text}\n"
            if images:
                content += f"\n[This slide contains {len(images)} related image(s)]\n"

            chunks.append(
                {
                    "content": content,
                    "filename": os.path.basename(file_path),
                    "filepath": file_path,
                    "filetype": ".pptx",
                    "page_number": slide_idx,
                    "images": images,
                }
            )

        return chunks

    def _load_docx(self, file_path: str) -> str:
        try:
            return docx2txt.process(file_path)
        except Exception as exc:
            print(f"[loader] docx read failed: {exc}")
            return ""

    def _load_txt(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as exc:
            print(f"[loader] txt read failed: {exc}")
            return ""

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------
    def _wrap_single_text_chunk(self, file_path: str, text: str) -> List[Dict]:
        if not text:
            return []
        return [
            {
                "content": text,
                "filename": os.path.basename(file_path),
                "filepath": file_path,
                "filetype": os.path.splitext(file_path)[1].lower(),
                "page_number": 0,
                "images": [],
            }
        ]

    def _extract_images_from_pdf_page(
        self, doc: fitz.Document, page_idx: int, pdf_name: str
    ) -> List[str]:
        """Save images from a single PDF page; returns file paths."""
        saved: List[str] = []
        page = doc[page_idx]
        for img_idx, img in enumerate(page.get_images(full=True), 1):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                ext = base_image["ext"]
                width = base_image["width"]
                height = base_image["height"]
                if width < 50 or height < 50:
                    continue
                fname = f"{pdf_name}_p{page_idx+1}_{img_idx}.{ext}"
                fpath = os.path.join(self.images_dir, fname)
                with open(fpath, "wb") as f:
                    f.write(image_bytes)
                saved.append(self._to_stored_image_path(fpath))
            except Exception as exc:
                print(f"[loader] pdf image extract failed: {exc}")
        return saved

    def _extract_images_from_pptx_slide(
        self, slide, slide_idx: int, pptx_name: str
    ) -> List[str]:
        saved: List[str] = []
        for shape_idx, shape in enumerate(slide.shapes, 1):
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                image = shape.image
                blob = image.blob
                if len(blob) < 1024:
                    continue
                ext = image.ext
                fname = f"{pptx_name}_slide{slide_idx}_{shape_idx}.{ext}"
                fpath = os.path.join(self.images_dir, fname)
                with open(fpath, "wb") as f:
                    f.write(blob)
                saved.append(self._to_stored_image_path(fpath))
            except Exception as exc:
                print(f"[loader] pptx image extract failed: {exc}")
        return saved

    def _to_stored_image_path(self, abs_path: str) -> str:
        rel = os.path.relpath(abs_path, start=self.project_root)
        return rel.replace("\\", "/")

    @staticmethod
    def encode_image_to_base64(image_path: str) -> str:
        """Utility to load a local image as base64 (for downstream vision prompts)."""
        if not os.path.exists(image_path):
            return ""
        with open(image_path, "rb") as img:
            return base64.b64encode(img.read()).decode("utf-8")
